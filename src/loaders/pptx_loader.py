"""PPTX loader — per-slide heading + shape text + tables + images + speaker notes."""
from __future__ import annotations

import logging
from pathlib import Path

from .elements import DocElement, elements_to_markdown, rows_to_markdown

log = logging.getLogger(__name__)


def load_pptx_elements(path: Path, *, extract_images: bool = False, image_out: Path | None = None) -> list[DocElement]:
    try:
        from pptx import Presentation
        from pptx.util import Emu  # noqa: F401
    except ImportError as e:
        raise RuntimeError("python-pptx not installed. Run: pip install python-pptx") from e

    prs = Presentation(str(path))
    elements: list[DocElement] = []

    for si, slide in enumerate(prs.slides, start=1):
        elements.append(DocElement(kind="heading", content=f"Slide {si}", meta={"level": 2}))
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = "".join(run.text for run in para.runs).strip()
                        if txt:
                            elements.append(DocElement(kind="text", content=txt, meta={"slide": si}))
                if shape.has_table:
                    rows = [[(c.text or "").strip() for c in row.cells] for row in shape.table.rows]
                    md = rows_to_markdown(rows)
                    if md:
                        elements.append(DocElement(kind="table", content=md, meta={"slide": si}))
                if extract_images and image_out is not None and shape.shape_type == 13:  # PICTURE
                    try:
                        img = shape.image
                        ext = img.ext or "png"
                        image_out.mkdir(parents=True, exist_ok=True)
                        tgt = image_out / f"{path.stem}-s{si}-{shape.shape_id}.{ext}"
                        tgt.write_bytes(img.blob)
                        elements.append(DocElement(kind="image", content="", meta={"slide": si, "path": str(tgt)}))
                    except Exception:
                        continue
            except Exception as e:
                log.debug("pptx shape skipped", extra={"metadata": {"error": str(e)[:120]}})
        # Speaker notes
        notes = slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else ""
        if notes:
            elements.append(DocElement(kind="text", content=f"_Speaker notes:_ {notes}", meta={"slide": si, "notes": True}))

    return elements


def load_pptx(path: Path) -> str:
    return elements_to_markdown(load_pptx_elements(path))
